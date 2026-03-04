import io
import fitz  # PyMuPDF for reliable PDF parsing
from docx import Document
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_community.utilities import GoogleSerperAPIWrapper
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.memory import MemorySaver
import streamlit as st
import re
from fpdf import FPDF
from gtts import gTTS

# Page Config for Better UI
st.set_page_config(page_title="StudyGenie AI", page_icon="🤖", layout="wide")

load_dotenv()

# Custom CSS for Sidebar Look
st.markdown("""
    <style>
    /* Distinct Sidebar Style */
    section[data-testid="stSidebar"] {
        background-color: #1e293b !important;
        border-right: 1px solid #334155;
    }
    
    /* Style buttons in the sidebar */
    .stSidebar .stButton > button {
        background: linear-gradient(135deg, #6366f1 0%, #a855f7 100%);
        color: white;
        border: none;
        border-radius: 8px;
        transition: all 0.3s ease;
    }
    </style>
""", unsafe_allow_html=True)

# Memory and Session State Initialization
def reset_chat():
    st.session_state.memory = MemorySaver()
    st.session_state.history = []
    st.session_state.materials_text = ""
    st.rerun()

if "memory" not in st.session_state:
    st.session_state.memory = MemorySaver()
    st.session_state.history = []
    st.session_state.materials_text = ""

# Sidebar Configuration
st.sidebar.title("🤖 Chat Controls")

# New Chat Button (ChatGPT style)
if st.sidebar.button("➕ New Chat", use_container_width=True):
    reset_chat()

# Initialize LLM and Tools
llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
search = GoogleSerperAPIWrapper()
tools = [search.run]

# Agent Setup with Study Helper Prompt
system_prompt = (
    "You are StudyGenie AI, a dedicated and supportive study assistant. "
    "Your mission is to help students understand their study materials, answer questions clearly, "
    "summarize content, and provide detailed explanations. "
    "\n\nCRITICAL RULE: When study materials are provided, you MUST prioritize information from those materials. "
    "If the answer is found in the uploaded documents, use it as the primary source. "
    "If the answer isn't in the material, you can use your general knowledge or search the web if needed, "
    "but clearly specify when you are doing so. "
    "\n\nAlways maintain an encouraging, academic, and professional tone."
)

agent = create_react_agent(
    model=llm,
    tools=tools,
    checkpointer=st.session_state.memory
)

# Web Interface Header
st.title("StudyGenie AI 🤖")
st.markdown("### Your Intelligent Study Companion")
st.write("Upload your notes, books, or data and ask anything related to your studies!")

# Function to create PDF from text
def create_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    
    # Basic cleanup for characters not in latin-1 (default for helvetica)
    # Often AI responses have emojis or special quotes.
    # Replace common non-latin-1 characters to avoid errors
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2018', "'").replace('\u2019', "'").replace('\u201c', '"').replace('\u201d', '"')
    safe_text = text.encode('latin-1', 'ignore').decode('latin-1')
    
    pdf.multi_cell(0, 10, txt=safe_text)
    return bytes(pdf.output())

# Function to convert text to speech
def read_aloud(text):
    # Clean text from markdown for better speech
    clean_text = re.sub(r'[*_#`]', '', text)
    tts = gTTS(text=clean_text, lang='en')
    fp = io.BytesIO()
    tts.write_to_fp(fp)
    fp.seek(0)
    return fp

def extract_text(files, format_choice):
    text_content = ""
    for file in files:
        if format_choice == "Text":
            try:
                text_content += file.read().decode("utf-8") + "\n"
            except Exception:
                file.seek(0)
                text_content += file.read().decode("latin-1", errors="replace") + "\n"
        elif format_choice == "PDF":
            # Using fitz for much better text extraction especially on large files
            doc = fitz.open(stream=file.read(), filetype="pdf")
            for page in doc:
                text_content += page.get_text() + "\n"
            doc.close()
        elif format_choice == "Word Document":
            doc = Document(file)
            for para in doc.paragraphs:
                text_content += para.text + "\n"
        elif format_choice == "PPT":
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif format_choice == "Excel" or format_choice == "CSV":
            df = pd.read_excel(file) if format_choice == "Excel" else pd.read_csv(file)
            text_content += df.to_string() + "\n"
    return text_content

# Main Chat Display
for idx, message in enumerate(st.session_state.history):
    st.chat_message(message["role"]).markdown(message["content"])
    
    # Handle export and read options for AI responses
    if message["role"] == "ai":
        col1, col2 = st.columns([1, 4])
        with col1:
            st.download_button(
                label="📥 PDF",
                data=create_pdf(message["content"]),
                file_name=f"StudyGenie_Response_{idx}.pdf",
                mime="application/pdf",
                key=f"export_chat_{idx}"
            )
        with col2:
            if st.button(f"🔊 Read Aloud", key=f"read_{idx}"):
                audio_fp = read_aloud(message["content"])
                st.audio(audio_fp, format="audio/mp3")

# Sidebar for Study Material Hub and Chat History
st.sidebar.title("📚 Study Material Hub")
choice = st.sidebar.selectbox(
    "Select Format",
    ("Text", "PDF", "Word Document", "PPT", "Excel", "CSV")
)

with st.sidebar:
    uploaded_files = st.file_uploader(f"Upload {choice}", type=None, accept_multiple_files=True)
    if uploaded_files:
        with st.spinner("Processing study materials..."):
            extracted_text = extract_text(uploaded_files, choice)
            if not extracted_text.strip():
                st.warning("No readable text found in the material. It might be a scanned document or an image-based PDF.")
            st.session_state.materials_text = extracted_text
            st.success(f"Successfully loaded {len(uploaded_files)} file(s)!")

    # Chat History Sidebar Section
    st.sidebar.divider()
    st.sidebar.title("📜 Recent Chat History")
    if not st.session_state.history:
        st.sidebar.info("No chat history yet.")
    else:
        for idx, msg in enumerate(st.session_state.history):
            if msg["role"] == "user":
                # Display questions as expandable items in sidebar
                with st.sidebar.expander(f"Q: {msg['content'][:30]}..."):
                    st.write(msg["content"])
                    # Find the AI response that follows
                    if idx + 1 < len(st.session_state.history):
                        st.info(st.session_state.history[idx+1]["content"])


# Chat input MUST be at the top level
query = st.chat_input("Ask a study question...")

if query:
    st.chat_message("user").markdown(query)
    st.session_state.history.append({"role": "user", "content": query})

    # Construct the query with context if available
    full_prompt = query
    if st.session_state.materials_text:
        # Increase context window to 15000 characters for better book/document coverage
        full_prompt = f"STUDY MATERIAL (Extracted Text):\n{st.session_state.materials_text[:15000]}...\n\nQUESTION: {query}"

    with st.spinner("Thinking..."):
        config = {"configurable": {"thread_id": "1"}}
        try:
            response = agent.invoke(
                {"messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": full_prompt}
                ]},
                config
            )
            
            # Handle various response formats from the model
            raw_answer = response["messages"][-1].content
            if isinstance(raw_answer, list):
                answer = "".join([part.get("text", "") if isinstance(part, dict) else str(part) for part in raw_answer])
            else:
                answer = str(raw_answer)

            st.session_state.history.append({"role": "ai", "content": answer})
            st.rerun()
            
        except Exception as e:
            if "RESOURCE_EXHAUSTED" in str(e) or "429" in str(e):
                st.error("🚫 **Quota Exceeded!** The daily limit for the Gemini API has been reached. Please try again later.")
            else:
                st.error(f"An error occurred: {str(e)}")
